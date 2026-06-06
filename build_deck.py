"""
Build the 10-slide deck for AtliQ Media.
RPC — Decoding the 2026 Tamil Nadu Assembly Election.

Run from the repo root:
    python build_deck.py

Expects the 4 chart PNGs at:
    charts/01_sankey.png
    charts/02_vote_share.png
    charts/03_margin.png
    charts/04_regional.png

If any image is missing, the script substitutes a labelled placeholder
rectangle so you can drop the screenshot in later in PowerPoint.

Outputs:
    deck/atliq_tn_election_2026.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

REPO = Path(__file__).resolve().parent
CHARTS = REPO / "charts"
DECK = REPO / "deck"
DECK.mkdir(exist_ok=True)
OUT = DECK / "atliq_tn_election_2026.pptx"

NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xE6, 0x39, 0x46)
GREY = RGBColor(0x36, 0x45, 0x4F)
MUTED = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
BORDER = RGBColor(0xC0, 0xC8, 0xD0)

F_HEAD = "Georgia"
F_BODY = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, *, font=F_BODY, size=14,
             color=GREY, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
    return tb


def add_rect(slide, x, y, w, h, fill, *, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def dark_bg(slide):
    add_rect(slide, 0, 0, 13.333, 7.5, NAVY)


def add_image_or_placeholder(slide, x, y, w, h, image_path, label):
    p = Path(image_path)
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
    else:
        add_rect(slide, x, y, w, h, LIGHT_GREY, line=BORDER)
        add_text(slide, x, y, w, h,
                 f"[ Insert: {label} ]\n(replace this box with your screenshot)",
                 font=F_BODY, size=14, color=MUTED, italic=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, page_no, dark=False):
    color = ICE if dark else MUTED
    add_text(slide, 0.5, 7.05, 8, 0.3,
             "Source: Election Commission of India  |  Non-partisan editorial analysis  |  AtliQ Media",
             size=9, color=color, italic=True)
    add_text(slide, 11.8, 7.05, 1.0, 0.3,
             f"{page_no} / 10", size=9, color=color,
             align=PP_ALIGN.RIGHT, italic=True)


# SLIDE 1 — COVER
s1 = prs.slides.add_slide(BLANK)
dark_bg(s1)
add_rect(s1, 0.7, 3.1, 0.18, 0.7, ACCENT)
add_text(s1, 1.0, 3.0, 11, 0.5,
         "RPC — An editorial deck for AtliQ Media",
         font=F_BODY, size=14, color=ICE, italic=True)
add_text(s1, 1.0, 3.55, 11.5, 1.5,
         "Decoding the 2026 Tamil Nadu\nAssembly Election",
         font=F_HEAD, size=44, color=WHITE, bold=True)
add_text(s1, 1.0, 5.5, 11.5, 0.6,
         "How 234 seats moved — told with only ECI data.",
         font=F_BODY, size=20, color=ICE)
add_text(s1, 1.0, 6.9, 11.5, 0.4,
         "Prepared by SK   |   Data: Election Commission of India   |   2026",
         font=F_BODY, size=10, color=ICE, italic=True)

# SLIDE 2 — HEADLINE
s2 = prs.slides.add_slide(BLANK)
dark_bg(s2)
add_text(s2, 0.5, 1.1, 12.3, 0.5, "THE HEADLINE",
         font=F_BODY, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(s2, 0.5, 1.7, 12.3, 2.6, "108 / 234",
         font=F_HEAD, size=160, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s2, 0.5, 4.5, 12.3, 0.6,
         "seats won by TVK in the 2026 Tamil Nadu Assembly election.",
         font=F_BODY, size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s2, 0.5, 5.2, 12.3, 0.6,
         "A party that did not exist in 2021.",
         font=F_BODY, size=20, color=ICE, italic=True, align=PP_ALIGN.CENTER)
add_footer(s2, 2, dark=True)

# SLIDE 3 — AGENDA
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, 13.333, 7.5, WHITE)
add_rect(s3, 0.5, 0.6, 0.15, 0.5, ACCENT)
add_text(s3, 0.75, 0.5, 12, 0.6, "Three stories. One thread.",
         font=F_HEAD, size=32, color=NAVY, bold=True)
add_text(s3, 0.75, 1.2, 12, 0.5,
         "What AtliQ's 60-minute show should be built around.",
         font=F_BODY, size=16, color=MUTED, italic=True)

def agenda_col(x, n, title, body):
    add_rect(s3, x, 2.2, 3.85, 4.0, LIGHT_GREY, line=BORDER)
    add_text(s3, x + 0.3, 2.4, 3.55, 0.7, n,
             font=F_HEAD, size=48, color=ACCENT, bold=True)
    add_text(s3, x + 0.3, 3.3, 3.55, 0.6, title,
             font=F_HEAD, size=20, color=NAVY, bold=True)
    add_text(s3, x + 0.3, 4.0, 3.55, 2.1, body,
             font=F_BODY, size=14, color=GREY)

agenda_col(0.75, "01", "The Flip",
           "234 seats. 163 changed hands. Where they moved — at a glance.")
agenda_col(4.75, "02", "The Evidence",
           "Where TVK's 34.9% statewide vote share came from — and what shrunk to make room.")
agenda_col(8.75, "03", "The Context",
           "Wins got narrower across the board. The pattern that connects every contest.")
add_footer(s3, 3)

# SLIDE 4 — SETUP
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, 13.333, 7.5, WHITE)
add_rect(s4, 0.5, 0.6, 0.15, 0.5, ACCENT)
add_text(s4, 0.75, 0.5, 12, 0.6,
         "Before the three stories — one frame.",
         font=F_HEAD, size=30, color=NAVY, bold=True)
add_text(s4, 0.75, 2.0, 6.5, 4.5,
         ("In a single election cycle, the seat map of Tamil Nadu shifted "
          "in ways the state has not seen in any single cycle for decades.\n\n"
          "All three of the stories in this deck flow from one structural "
          "change: a third pillar entered a contest the state had organised "
          "as a two-front race since 1967.\n\n"
          "The numbers do not require interpretation. They require attention."),
         font=F_BODY, size=17, color=GREY)
add_rect(s4, 8.0, 2.0, 4.8, 4.0, NAVY)
add_text(s4, 8.2, 2.3, 4.5, 0.5, "STAT THAT FRAMES IT",
         font=F_BODY, size=12, color=ACCENT, bold=True)
add_text(s4, 8.2, 2.85, 4.5, 1.6, "163 / 234",
         font=F_HEAD, size=66, color=WHITE, bold=True)
add_text(s4, 8.2, 4.8, 4.5, 0.5,
         "seats changed winning party.",
         font=F_BODY, size=16, color=ICE)
add_text(s4, 8.2, 5.3, 4.5, 0.5,
         "= 69.7% of every constituency in TN.",
         font=F_BODY, size=14, color=ICE, italic=True)
add_footer(s4, 4)

# SLIDE 5 — FLIP
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, 13.333, 7.5, WHITE)
add_text(s5, 0.5, 0.4, 3, 0.4, "01  THE FLIP",
         font=F_BODY, size=12, color=ACCENT, bold=True)
add_text(s5, 0.5, 0.8, 12.3, 0.7,
         "234 seats. Where they moved. 2021 → 2026.",
         font=F_HEAD, size=26, color=NAVY, bold=True)
add_image_or_placeholder(s5, 0.5, 1.7, 12.3, 4.6,
                         CHARTS / "01_sankey.png",
                         "Chart 1 — Sankey diagram (2021 → 2026 seat flows)")
add_text(s5, 0.5, 6.4, 12.3, 0.5,
         "Of TVK's 108 wins: 65 from DMK · 26 from AIADMK · 11 from INC · 6 from BJP/PMK/VCK/CPI(M)/OTHERS",
         font=F_BODY, size=13, color=GREY, italic=True, align=PP_ALIGN.CENTER)
add_footer(s5, 5)

# SLIDE 6 — VOTE SHARE
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, 13.333, 7.5, WHITE)
add_text(s6, 0.5, 0.4, 3, 0.4, "02  THE EVIDENCE",
         font=F_BODY, size=12, color=ACCENT, bold=True)
add_text(s6, 0.5, 0.8, 12.3, 0.7,
         "Where TVK's 34.9% statewide share came from.",
         font=F_HEAD, size=26, color=NAVY, bold=True)
add_image_or_placeholder(s6, 0.5, 1.7, 8.5, 4.8,
                         CHARTS / "02_vote_share.png",
                         "Chart 2 — Statewide vote share, 2021 vs 2026")
add_rect(s6, 9.4, 1.7, 3.5, 4.8, LIGHT_GREY, line=BORDER)
add_text(s6, 9.6, 1.9, 3.1, 0.4, "DROPS FROM 2021",
         font=F_BODY, size=11, color=ACCENT, bold=True)
add_text(s6, 9.6, 2.4, 3.1, 2.5,
         "DMK:     -13.5 pp\nAIADMK: -12.1 pp\nNTK:      -2.6 pp\nOTHERS:  -2.6 pp\nPMK:      -1.6 pp\nAMMK:     -1.5 pp\nINC:        -0.9 pp",
         font=F_BODY, size=15, color=GREY)
add_text(s6, 9.6, 5.0, 3.1, 0.4, "TVK GAIN",
         font=F_BODY, size=11, color=ACCENT, bold=True)
add_text(s6, 9.6, 5.4, 3.1, 0.7, "+34.9 pp",
         font=F_HEAD, size=28, color=NAVY, bold=True)
add_text(s6, 9.6, 6.0, 3.1, 0.4, "= sum of all the drops.",
         font=F_BODY, size=12, color=GREY, italic=True)
add_footer(s6, 6)

# SLIDE 7 — MARGIN
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, 13.333, 7.5, WHITE)
add_text(s7, 0.5, 0.4, 3, 0.4, "03  THE CONTEXT",
         font=F_BODY, size=12, color=ACCENT, bold=True)
add_text(s7, 0.5, 0.8, 12.3, 0.7,
         "Wins got narrower across the board.",
         font=F_HEAD, size=26, color=NAVY, bold=True)
add_image_or_placeholder(s7, 0.5, 1.7, 8.5, 4.8,
                         CHARTS / "03_margin.png",
                         "Chart 3 — Vote-share buckets, 2021 vs 2026")

def stat_card(y, label, val):
    add_rect(s7, 9.4, y, 3.5, 1.45, LIGHT_GREY, line=BORDER)
    add_text(s7, 9.6, y + 0.15, 3.1, 0.4, label,
             font=F_BODY, size=11, color=ACCENT, bold=True)
    add_text(s7, 9.6, y + 0.55, 3.1, 0.8, val,
             font=F_HEAD, size=20, color=NAVY, bold=True)

stat_card(1.7, "DECISIVE WINS (>50%)", "70 -> 13")
stat_card(3.3, "MINORITY WINS (<35%)", "2 -> 64")
stat_card(4.9, "MEDIAN MARGIN", "9.5% -> 5.7%")
add_footer(s7, 7)

# SLIDE 8 — REGIONAL
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, 13.333, 7.5, WHITE)
add_text(s8, 0.5, 0.4, 8, 0.4, "REGIONAL VIEW",
         font=F_BODY, size=12, color=ACCENT, bold=True)
add_text(s8, 0.5, 0.8, 12.3, 0.7,
         "The pattern held across all six regions.",
         font=F_HEAD, size=26, color=NAVY, bold=True)
add_image_or_placeholder(s8, 0.5, 1.7, 12.3, 4.6,
                         CHARTS / "04_regional.png",
                         "Chart 4 — Seats by region, 2021 vs 2026")
add_text(s8, 0.5, 6.4, 12.3, 0.5,
         "TVK won seats in every region. Chennai Metro near-reversal: DMK 29->2  ·  TVK 0->29 of 32 seats.",
         font=F_BODY, size=13, color=GREY, italic=True, align=PP_ALIGN.CENTER)
add_footer(s8, 8)

# SLIDE 9 — EDITORIAL REC
s9 = prs.slides.add_slide(BLANK)
add_rect(s9, 0, 0, 13.333, 7.5, WHITE)
add_rect(s9, 0.5, 0.6, 0.15, 0.5, ACCENT)
add_text(s9, 0.75, 0.5, 12, 0.6,
         "Editorial recommendation",
         font=F_HEAD, size=30, color=NAVY, bold=True)
add_text(s9, 0.75, 1.15, 12, 0.5,
         "How the 60-minute show opens, pivots, and closes.",
         font=F_BODY, size=15, color=MUTED, italic=True)

def beat(y, label, title, body):
    add_text(s9, 0.75, y, 0.6, 0.6, label,
             font=F_HEAD, size=28, color=ACCENT, bold=True)
    add_text(s9, 1.5, y, 6.0, 0.5, title,
             font=F_HEAD, size=17, color=NAVY, bold=True)
    add_text(s9, 1.5, y + 0.5, 6.0, 1.0, body,
             font=F_BODY, size=13, color=GREY)

beat(2.0, "1", "OPEN with one number",
     "Anchor opens with: \"108 of 234.\" Pause. Then: \"A party that did not exist in 2021.\"")
beat(3.6, "2", "PIVOT to the flip",
     "Show the Sankey. Voiceover: \"And here is where those 108 came from.\" 65 DMK + 26 AIADMK + 11 INC + 6 others.")
beat(5.2, "3", "CLOSE with the margin",
     "Show the bucket chart. \"Decisive wins fell from 70 to 13. Minority wins rose from 2 to 64. This is what a tighter contest looks like.\"")

add_rect(s9, 8.5, 2.0, 4.4, 4.5, NAVY)
add_text(s9, 8.7, 2.15, 4.0, 0.4, "TONE DISCIPLINE",
         font=F_BODY, size=12, color=ACCENT, bold=True)
add_text(s9, 8.7, 2.6, 4.0, 3.7,
         ("Every line must read identically to a supporter of any party.\n\n"
          "WORDS TO AVOID:\n"
          "    smashed - collapsed - swept\n"
          "    destroyed - ended - broke\n\n"
          "WORDS TO KEEP:\n"
          "    won - lost - moved - shifted\n"
          "    pulled from - came from\n\n"
          "NEVER attribute causation. The data shows what moved -- not why."),
         font=F_BODY, size=12, color=ICE)
add_footer(s9, 9)

# SLIDE 10 — LIMITATIONS
s10 = prs.slides.add_slide(BLANK)
dark_bg(s10)
add_text(s10, 0.5, 0.5, 12.3, 0.5,
         "WHAT THIS DATA CAN AND CANNOT SAY",
         font=F_BODY, size=12, color=ACCENT, bold=True)
add_text(s10, 0.5, 1.0, 12.3, 0.8,
         "Limitations, methodology, and honesty.",
         font=F_HEAD, size=28, color=WHITE, bold=True)

def lim_item(x, y, w, head, body):
    add_text(s10, x, y, w, 0.4, head,
             font=F_BODY, size=12, color=ACCENT, bold=True)
    add_text(s10, x, y + 0.45, w, 1.4, body,
             font=F_BODY, size=13, color=ICE)

lim_item(0.5, 2.1, 5.8, "TURNOUT (2026)",
         "The 2026 turnout column is blank in the source CSV. "
         "Turnout-driven stories were not analysed in this deck.")
lim_item(7.0, 2.1, 5.8, "FORM-20 NOT YET RELEASED",
         "All 2026 numbers are from the ECI live results portal "
         "ahead of the final audited Form-20.")
lim_item(0.5, 3.9, 5.8, "DUPLICATE CANDIDATES",
         "39 same-name cases retained. Pattern is the known dummy-candidate "
         "ballot-confusion tactic plus common-name collisions among independents.")
lim_item(7.0, 3.9, 5.8, "PARTY NORMALISATION",
         "13 major parties tracked explicitly; ~110 fringe parties bucketed as OTHERS. "
         "Does not affect headline numbers.")
lim_item(0.5, 5.7, 5.8, "WHAT WE DO NOT CLAIM",
         "No causal explanation for any flip. No prediction of future elections. "
         "No assessment of any party, leader, region or community.")
lim_item(7.0, 5.7, 5.8, "SOURCES & REPRODUCIBILITY",
         "Code + notebook + data: github.com/<SK>/tn-election-2026-decoded\n"
         "Primary source: ECI -- Cross-check: Trivedi Centre, Ashoka University.")
add_footer(s10, 10, dark=True)

prs.save(OUT)
print(f"Deck written: {OUT}")
print(f"Slides: {len(prs.slides)}")
